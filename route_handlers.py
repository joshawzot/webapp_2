# route_handlers.py

from my_flask_app import app, cache, redis_client
from db_operations import *
from config import *
from utilities import *
# my_flask_app.py
from flask import Flask, request, jsonify, render_template, redirect
from flask_sqlalchemy import SQLAlchemy
from flask_mail import Mail, Message
from flask_apscheduler import APScheduler
from flask_caching import Cache
from flask_talisman import Talisman  # Ensure Talisman is used if security headers are required

import redis  # Used for Redis client initialization
import json  # Used for JSON handling, if necessary
import uuid  # Used for generating unique identifiers
import random  # Used if random operations are needed

from config import DB_CONFIG, APP_SECRET_KEY  # Import configuration constants

# Standard library imports
import os, hashlib, base64, json, uuid, time
from io import BytesIO
from tempfile import TemporaryFile, NamedTemporaryFile, TemporaryDirectory
from urllib.parse import quote_plus

# External libraries
import pandas as pd
import requests
import boto3
import mysql.connector
from flask import Flask, request, jsonify, make_response, redirect, url_for, session, flash, send_file, render_template, render_template_string
from werkzeug.security import check_password_hash, generate_password_hash
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
import zipfile
import numpy as np
import csv
import io
from PIL import Image

# Custom module imports
from generate_plot_endurance import generate_plot_endurance
from generate_plot import generate_plot
from generate_plot_normal_combined import generate_plot_normal_combined
from generate_plot_combined import generate_plot_combined
from generate_plot_separate import generate_plot_separate

from flask_caching import Cache
cache = Cache(app, config={'CACHE_TYPE': 'simple'})

@app.route('/create-db')
def create_db_page():
    conn = create_connection()
    cursor = conn.cursor()
    databases = get_all_databases(cursor)  # Fetch all database names
    cursor.close()
    conn.close()

    return render_template('create_db_page.html')

@app.route('/save-txt-content/<database>/<table_name>', methods=['POST'])
def save_txt_content(database, table_name):
    try:
        content = request.json['content']
        connection = create_connection(database)
        cursor = connection.cursor()
        query = f"UPDATE `{table_name}` SET content = %s WHERE content IS NOT NULL LIMIT 1"
        cursor.execute(query, (content,))
        connection.commit()
        close_connection()
        return "Content saved successfully", 200
    except mysql.connector.Error as err:
        return str(err), 400

@app.route('/')  # Defines a route for the root URL
def home():
    print("home")
    return redirect(url_for('home_page'))

@app.route('/home-page', methods=['GET'])
def home_page():
    """Get a list of databases available."""
    try:
        conn = create_connection()
        cursor = conn.cursor()
        databases = get_all_databases(cursor)  # Fetch all database names
        cursor.close()
        conn.close()

        return render_template('home_page.html', databases=databases)

    except mysql.connector.Error as err:
        return str(err)

@app.route('/list-tables', methods=['POST', 'GET'])
def list_tables():
    if request.method == 'POST':
        session['database'] = request.form.get('database')

    database = session.get('database')

    tables = fetch_tables(database)  # Retrieve table data from the database
    table_names = ','.join(table['table_name'] for table in tables)
    print("table_names:", table_names)

    plot_function = "None"  # This could also be dynamically set based on POST or other conditions

    return render_template('list_tables.html', tables=tables, table_names=table_names, database=database, plot_function=plot_function)

@app.route('/view-table/<database>/<table_name>', methods=['GET'])
def view_table(database, table_name):
    """View the content of a specific table."""
    print('database:', database)
    print('table_name:', table_name)
    try:
        connection = create_connection(database)
        cursor = connection.cursor()
        if table_name.endswith('_txt'):
            query = f"SELECT content FROM `{table_name}` LIMIT 1"
            results = fetch_data(cursor, query)
            close_connection()
            content = results[0][0] if results else ''
            print('------')
            print(table_name)
            return render_template('table_txt.html', database=database, content=content, table_name=table_name)
        else:
            query = f"SELECT * FROM `{table_name}`"
            results = fetch_data(cursor, query)
            column_names = [desc[0] for desc in cursor.description]
            close_connection()
            return render_template('table.html', results=results, column_names=column_names)
    except mysql.connector.Error as err:
        return str(err)

# Define a dictionary to map database names to functions
generate_plot_functions = {
    "generate_plot_endurance": generate_plot_endurance,
    "generate_plot": generate_plot,
    "generate_plot_normal_combined": generate_plot_normal_combined,
    "generate_plot_combined": generate_plot_combined,
    "generate_plot_separate": generate_plot_separate,
}
'''
@app.route('/render-plot/<unique_id>')
def render_plot(unique_id):
    data_type = "avg_std"
    #return download_csv2(unique_id, data_type)

    # Regular flow to fetch and render plot data
    cache_key = f"plot_data_{unique_id}"
    cached_plot_data = cache.get(cache_key)

    if cached_plot_data:
        return render_template('plot.html', plot_data=cached_plot_data)

    stored_data_json = redis_client.get(unique_id)
    if not stored_data_json:
        return "Error: Invalid ID or Data Expired", 404

    stored_data = json.loads(stored_data_json)
    database = stored_data["database"]
    table_names = stored_data["table_name"].split(',')
    form_data = stored_data["form_data"]
    plot_function = stored_data["plot_function"]

    generate_plot_function = generate_plot_functions.get(plot_function)
    if generate_plot_function is None:
        return "Error: Invalid plot function selection", 400

    try:
        plot_data = generate_plot_function(table_names, database, form_data)
        cache.set(cache_key, plot_data, timeout=None)
        return render_template('plot.html', plot_data=plot_data)
    except Exception as e:
        return f"Error: {e}", 500
'''
'''
@app.route('/render-plot/<database>/<table_name>/<plot_function>')
def render_plot(database, table_name, plot_function):
    data_type = "avg_std"
    form_data_json = request.args.get('form_data')
    
    if not form_data_json:
        return "Error: Form data not provided", 400

    try:
        form_data = json.loads(form_data_json)
    except json.JSONDecodeError:
        return "Error: Invalid form data", 400

    generate_plot_function = generate_plot_functions.get(plot_function)
    if generate_plot_function is None:
        return "Error: Invalid plot function selection", 400

    try:
        table_names = table_name.split(',')
        plot_data = generate_plot_function(table_names, database, form_data)
        return render_template('plot.html', plot_data=plot_data)
    except Exception as e:
        return f"Error: {e}", 500
'''
@app.route('/render-plot/<database>/<table_name>/<plot_function>')
def render_plot(database, table_name, plot_function):
    form_data_json = request.args.get('form_data')
    
    if not form_data_json:
        return "Error: Form data not provided", 400

    try:
        form_data = json.loads(form_data_json)
    except json.JSONDecodeError:
        return "Error: Invalid form data", 400

    cache_key = f"{database}_{table_name}_{plot_function}_{form_data_json}"
    cached_plot_data = cache.get(cache_key)

    if cached_plot_data:
        return render_template('plot.html', plot_data=cached_plot_data)

    generate_plot_function = generate_plot_functions.get(plot_function)
    if generate_plot_function is None:
        return "Error: Invalid plot function selection", 400

    try:
        table_names = table_name.split(',')
        plot_data = generate_plot_function(table_names, database, form_data)
        cache.set(cache_key, plot_data, timeout=3600)  # Cache for 1 hour
        return render_template('plot.html', plot_data=plot_data)
    except Exception as e:
        return f"Error: {e}", 500

@app.route('/download_csv/<unique_id>/<data_type>')
def download_csv2(unique_id, data_type):
    # Retrieve plot data either from cache or Redis
    cache_key = f"plot_data_{unique_id}"
    plot_data = cache.get(cache_key)
    if not plot_data:
        stored_data_json = redis_client.get(unique_id)
        if not stored_data_json:
            return "Error: Data not found", 404
        stored_data = json.loads(stored_data_json)
        database = stored_data["database"]
        table_names = stored_data["table_name"].split(',')
        form_data = stored_data["form_data"]
        plot_function = stored_data["plot_function"]
        generate_plot_function = generate_plot_functions.get(plot_function)
        if not generate_plot_function:
            return "Error: Invalid plot function selection", 400
        plot_data = generate_plot_function(table_names, database, form_data)

    # Generate CSV based on plot_data and data_type
    if data_type == "avg_std":
        avg_values, std_values, table_names, selected_groups = plot_data
        header = ["State"] + [f"{table_name}" for table_name in table_names] + ["Row Avg", "Row Std Dev"]
        table_data = [header]
        column_data = [[] for _ in table_names]

        for i, group in enumerate(selected_groups):
            row = [f"State {group}"]
            row_data = []

            for j, table_avg in enumerate(avg_values):
                avg = table_avg[i]
                row.append(f"{avg:.2f}")
                row_data.append(avg)
                column_data[j].append(avg)

            row_avg = np.mean(row_data)
            row_std = np.std(row_data)
            row.extend([f"{row_avg:.2f}", f"{row_std:.2f}"])
            table_data.append(row)

        col_avgs = [np.mean(col) for col in column_data]
        col_stds = [np.std(col) for col in column_data]
        table_data.append(["Col Avg"] + [f"{avg:.2f}" for avg in col_avgs] + ["-", "-"])
        table_data.append(["Col Std Dev"] + [f"{std:.2f}" for std in col_stds] + ["-", "-"])
        return generate_csv_response(table_data, "avg_std_data.csv")

    elif data_type in ["sigma", "ppm", "us"]:
        ber_results, _ = plot_data
        headers = ["State/Transition"] + [name for name in table_names] + ["Row Avg"]
        data_collections = [headers[:], headers[:], headers[:]]

        grouped_data = {}
        for entry in ber_results:
            key = entry[1]
            if key not in grouped_data:
                grouped_data[key] = []
            grouped_data[key].append((entry[2], entry[3], entry[4]))

        for key, values in grouped_data.items():
            rows = [[key], [key], [key]]
            for val in values:
                rows[0].append(f"{val[0]:.4f}")
                rows[1].append(f"{int(val[1])}")
                rows[2].append(f"{int(val[2])}")
            for row in rows:
                avg = np.mean([float(v) for v in row[1:]])
                row.append(f"{avg:.4f}")

            data_collections[0].append(rows[0])
            data_collections[1].append(rows[1])
            data_collections[2].append(rows[2])

        index = {"sigma": 0, "ppm": 1, "us": 2}[data_type]
        filename = f"{data_type}_data.csv"
        return generate_csv_response(data_collections[index], filename)

def generate_csv_response(data, filename):
    csv_output = StringIO()
    for row in data:
        csv_output.write(','.join(str(item) for item in row) + '\n')
    csv_output.seek(0)
    response = make_response(csv_output.getvalue())
    response.headers["Content-Disposition"] = f"attachment; filename={filename}"
    response.headers["Content-type"] = "text/csv"
    return response

@app.route('/download_csv')
def download_csv():
    database = request.args.get('database')
    table_name = request.args.get('table_name')

    try:
        # Generate the CSV data
        csv_data = get_csv_from_table(database, table_name)
        if csv_data is None:
            return "Error generating CSV file", 500

        # Create a response with the CSV data as a downloadable file
        response = make_response(csv_data)
        response.headers['Content-Disposition'] = f'attachment; filename={table_name}.csv'
        response.mimetype = 'text/csv'
        return response
    except Exception as e:
        return str(e), 500

def fetch_all_numeric_data(table_name, database_name):
    # Establish a connection to the database
    connection = create_connection(database_name)
    data_dimension = "Unknown"
    try:
        with connection.cursor() as cursor:
            query = f"SELECT * FROM `{table_name}`"
            # Execute the query and fetch data
            cursor.execute(query)
            data = cursor.fetchall()
            # Determine the dimension of the data
            if data:
                # Check if the data has only one column (1D)
                if len(data[0]) == 1:
                    data_dimension = "1D"
                # If data has more than one column, consider it 2D
                else:
                    data_dimension = "2D"
                # If your application has a specific interpretation of 3D data,
                # you can add another condition here to set data_dimension to "3D"
    finally:
        connection.close()

    return data, data_dimension
'''
@app.route('/view-plot/<database>/<table_name>/<plot_function>', methods=['GET', 'POST'])
def view_plot(database, table_name, plot_function):
    print("view_plot")
    if request.method == "POST":
        print("POST:::::::::::::::::::::::::::::::::")
        # Check if the user has made a plot function choice
        plot_function_choice = request.form.get('plot_choice')
        if plot_function_choice:
            plot_function = plot_function_choice
            if plot_function in ["generate_plot", "generate_plot_combined", "generate_plot_separate", "generate_plot_normal_combined"]:
                return render_template(f'input_form_generate_plot.html', database=database, table_name=table_name, plot_function=plot_function)             
            else:
                return jsonify({"error": "choice not selected"}), 400

        # Ensure plot_function has a value before proceeding
        if plot_function:
            print(f"plot_function: {plot_function}")  # Now plot_function should have a value
            
            # Conditional logic to handle form data based on plot_function
            if plot_function in ["generate_plot", "generate_plot_combined", "generate_plot_separate", "generate_plot_normal_combined"]:
                form_data_handlers = {
                    "generate_plot": get_form_data_generate_plot,
                    "generate_plot_combined": get_form_data_generate_plot,
                    "generate_plot_separate": get_form_data_generate_plot,
                    "generate_plot_normal_combined": get_form_data_generate_plot,
                }
                form_data = form_data_handlers[plot_function](request.form)

                # Store the form data with a unique identifier in Redis
                unique_id = str(uuid.uuid4())
                redis_client.set(unique_id, json.dumps({
                    "database": database,
                    "table_name": table_name,
                    "plot_function": plot_function,
                    "form_data": form_data
                }))

                new_url = f"/render-plot/{unique_id}"
                return redirect(new_url)
            else:
                return jsonify({"error": "Invalid plot function selection"}), 400
        else:
            return jsonify({"error": "Plot function not selected"}), 400

    else:  # GET request handling
        #print('fuck')
        #print(len(table_names))
        table_names = table_name.split(',')
        print(table_names)
        print("GET:::::::::::::::::::::::::::::::::")
        return render_template('choose_plot_function_form.html')
'''

@app.route('/view-plot/<database>/<table_name>/<plot_function>', methods=['GET', 'POST'])
def view_plot(database, table_name, plot_function):
    print("view_plot")
    if request.method == "POST":
        print("POST:::::::::::::::::::::::::::::::::")
        plot_function_choice = request.form.get('plot_choice')
        if plot_function_choice:
            plot_function = plot_function_choice
            if plot_function in generate_plot_functions:
                return render_template('input_form_generate_plot.html', database=database, table_name=table_name, plot_function=plot_function)             
            else:
                return jsonify({"error": "Invalid plot function selection"}), 400

        if plot_function:
            print(f"plot_function: {plot_function}")
            if plot_function in generate_plot_functions:
                form_data = get_form_data_generate_plot(request.form)
                form_data_json = json.dumps(form_data)

                return redirect(f"/render-plot/{database}/{table_name}/{plot_function}?form_data={form_data_json}")
            else:
                return jsonify({"error": "Invalid plot function selection"}), 400
        else:
            return jsonify({"error": "Plot function not selected"}), 400
    else:
        table_names = table_name.split(',')
        print(table_names)
        print("GET:::::::::::::::::::::::::::::::::")
        return render_template('choose_plot_function_form.html')

@app.route('/upload-file', methods=['POST'])
def upload_file():
    print("upload_file()")
    if 'db_name' not in request.form:
        return jsonify(error="No database selected"), 400

    if 'files[]' not in request.files:
        print("No files part in request.files")  # Log missing files part
        return jsonify(error="No files part in the request"), 400

    files = [f for f in request.files.getlist('files[]') if f.filename]
    if not files:
        print("No files detected")  # Log no files detected
        return jsonify(error="No files detected"), 400

    for file in files:
        print(f"File: {file.filename}, MIME Type: {file.mimetype}")
    
    print(files)

    if not files:
        return jsonify(error="No files selected"), 400
    db_name = request.form['db_name']
    engine = create_db_engine(db_name)

    results = []
    for file in files:
        filename = sanitize_table_name(file.filename)
        file_extension = filename.rpartition('_')[-1]
        print("file_extension:", file_extension)
        file_stream = BytesIO(file.read())

        try:
            df = process_file(file_stream, file_extension, db_name)
            if df.shape[1] > 1017:
                df = df.transpose()
            if not df.empty:
                df.to_sql(filename, engine, if_exists='replace', index=False)
                results.append(f"{filename} uploaded successfully")
            else:
                results.append(f"No data to upload for {filename}. Dataframe is empty.")
        except Exception as e:
            error_msg = f"Error processing {filename}: {str(e)}"
            results.append(error_msg)

    return jsonify(results=results)

@app.route('/delete-record/<database>/<table_name>', methods=['DELETE'])  # delete a table
def delete_record(database, table_name):
    try:
        connection = create_connection(database)
        cursor = connection.cursor()

        query = f"DROP TABLE `{table_name}`"
        cursor.execute(query)

        connection.commit()
        close_connection()

        return "Record deleted successfully", 200
    except mysql.connector.Error as err:
        return str(err), 400

@app.route('/delete-records/<database>', methods=['DELETE'])  # delete multiple tables
def delete_records(database):
    try:
        tables = request.json['tables']
        connection = create_connection(database)
        cursor = connection.cursor()

        for table_name in tables:
            query = f"DROP TABLE `{table_name}`"
            cursor.execute(query)

        connection.commit()
        close_connection()

        return "Records deleted successfully", 200
    except mysql.connector.Error as err:
        return str(err), 400

@app.route('/create-database', methods=['POST'])
def create_database():
    db_name = request.form.get('newDatabaseName')  #"newDatabaseName" passed from create_db_page.html

    if not db_name:
        #return jsonify({"error": "No database name provided"}), 400
        return "No Folder name provided", 400
    
    else:
        create_db(db_name)
        #return jsonify({"message": f"Database {db_name} created successfully"}), 400
        return f"Folder {db_name} created successfully", 400

@app.route('/download_pptx', methods=['POST'])
def download_pptx():
    template_path = '/home/lenovoi7/Desktop/webapp_2/pptx_template/template.pptx'

    plots = request.json.get('plots', [])  # Retrieve the Base64 encoded images from the POST request

    prs = Presentation(template_path)  # Open the template PowerPoint file as the base for the new presentation
    
    # Retrieve the Base64 encoded images from the POST request
    plots = request.json.get('plots', [])
    
    for plot_data in plots:
        # Decode each Base64 image
        image_data = base64.b64decode(plot_data.split(",")[-1])
        # Open the image for analysis
        image = Image.open(BytesIO(image_data))
        
        # Choose a slide layout (6 is usually a blank slide)
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all shapes (including text boxes) from the slide
        for shape in slide.shapes:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # Get the image size
        img_width, img_height = image.size
        # Get the slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Calculate the scaling factor to maintain aspect ratio
        ratio = min(slide_width / img_width, slide_height / img_height)
        new_width = int(img_width * ratio)
        new_height = int(img_height * ratio)
        
        # Center the image
        left = int((slide_width - new_width) / 2)
        top = int((slide_height - new_height) / 2)
        
        # Convert the image data back to a BytesIO object
        img_io = BytesIO(image_data)
        # Add the image to the slide
        slide.shapes.add_picture(img_io, left, top, width=new_width, height=new_height)
    
    # Prepare the presentation to be sent in the response
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    # Set up the response with the correct headers
    response = make_response(pptx_io.getvalue())
    response.headers.set('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response.headers.set('Content-Disposition', 'attachment; filename="Downloaded_Presentation.pptx"')
    
    return response

@app.route('/rename-table/<database>/<old_name>/<new_name>', methods=['PUT'])
def rename_table(database, old_name, new_name):
    try:
        connection = create_connection(database)
        cursor = connection.cursor()

        # Construct and execute the SQL command
        sql = f"RENAME TABLE `{old_name}` TO `{new_name}`"
        cursor.execute(sql)
        connection.commit()
        close_connection()
        return "Table renamed successfully!", 200
    except mysql.connector.Error as err:
        print("Error occurred:", err)
        return str(err), 400

@app.route('/move-tables/<source_db>/<target_db>', methods=['PUT'])
def move_tables_route(source_db, target_db):
    try:
        success = move_tables(source_db, target_db)
        if success:
            return f"Tables moved from {source_db} to {target_db} successfully!", 200
        else:
            return f"Failed to move tables from {source_db} to {target_db}.", 400
    except Exception as err:
        print("Error occurred:", err)
        return str(err), 500

@app.route('/copy-all-tables', methods=['POST'])
def copy_all_tables_route():
    source_db = request.form.get('sourceDatabase')
    target_db = request.form.get('targetDatabase')

    try:
        copy_all_tables(source_db, target_db)
        return jsonify({"message": f"All tables copied from {source_db} to {target_db} successfully"}), 200
    except Exception as err:
        return jsonify({"error": str(err)}), 500

@app.route('/moveSelectedTables', methods=['POST'])
def move_selected_tables():
    data = request.json
    print("Received data for moving tables:", data)
    source_db = data['sourceDb']
    target_db = data['targetDb']
    table_names = data['tableNames']
    response = {"success": [], "failed": []}

    for table_name in table_names:
        try:
            if move_tables_2(source_db, target_db, table_name):
                response["success"].append(table_name)
            else:
                response["failed"].append(table_name)
        except Exception as e:
            print(f"Failed to move {table_name}: {e}")
            response["failed"].append(table_name)

    if response["failed"]:
        return jsonify({"success": False, "message": "Some tables failed to move.", "details": response}), 400
    return jsonify({"success": True, "message": "All selected tables moved successfully."})

@app.route('/copySelectedTables', methods=['POST'])
def copy_selected_tables():
    data = request.json
    print("Received data for copying tables:", data)
    source_db = data['sourceDb']
    target_db = data['targetDb']
    table_names = data['tableNames']
    response = {"success": [], "failed": []}
    print("Attempting to copy tables:", table_names)

    for table_name in table_names:
        try:
            if copy_tables_2(source_db, target_db, table_name):
                response["success"].append(table_name)
            else:
                response["failed"].append(table_name)
        except Exception as e:
            print(f"Failed to copy {table_name}: {e}")
            response["failed"].append(table_name)

    if response["failed"]:
        return jsonify({"success": False, "message": "Some tables failed to copy.", "details": response}), 400
    return jsonify({"success": True, "message": "All selected tables copied successfully."})

@app.route('/getDatabases', methods=['GET'])
def get_databases():
    # Create a database connection
    try:
        conn = create_connection()
        cursor = conn.cursor()
        
        # Use your function to get database names
        databases = get_all_databases(cursor)
        
        # Don't forget to close the cursor and connection when done
        cursor.close()
        conn.close()
        
        # Return the list of databases as a JSON response
        return jsonify(databases)

    except mysql.connector.Error as err:
        # In case of any database connection errors, return an error message
        return jsonify({"error": str(err)}), 500

@app.route('/add-item', methods=['POST'])
def create_item():
    item_id = request.form['item_id']
    item_data = request.form['item_data']
    try:
        response = table.put_item(
            Item={
                'ID': item_id,
                'item_data': item_data
            }
        )
        return redirect(url_for('home'))
    except ClientError as e:
        return jsonify({'error': str(e)}), 500

@app.route('/items', methods=['GET'])
def list_items():
    try:
        response = table.scan()
        items = response['Items']
        return jsonify(items)
    except ClientError as e:
        return jsonify({'error': str(e)}), 500
